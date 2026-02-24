"""
Test script to create a sample PPTX and test the converter fix.
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt

# Create a simple PowerPoint presentation
print("Creating test PowerPoint file...")
prs = Presentation()

# Add a title slide
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Test Presentation"
subtitle.text = "Sample file for testing PDF conversion"

# Add a content slide
bullet_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = "Content Slide"
tf = body_shape.text_frame
tf.text = "First point"
p = tf.add_paragraph()
p.text = "Second point"
p.level = 0
p = tf.add_paragraph()
p.text = "Third point"
p.level = 1

# Save the test file
test_file = Path("input/pp/test_sample.pptx")
test_file.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(test_file))
print(f"Test file created: {test_file}")

# Now test the conversion
print("\nTesting conversion...")
from src.core.powerpoint_converter import PowerPointConverter
from src.config import PDFConversionSettings

converter = PowerPointConverter()
output_file = Path("output/test_sample.pdf")
output_file.parent.mkdir(parents=True, exist_ok=True)

try:
    result = converter.convert(test_file, output_file)
    print(f"✓ Conversion successful: {result}")
except Exception as e:
    print(f"✗ Conversion failed: {e}")
    import traceback
    traceback.print_exc()
